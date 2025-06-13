from pptx import Presentation
from pptx.util import Pt
import pandas as pd
import os
import comtypes.client

def pptx_to_pdf(pptx_file, pdf_file):
    if os.path.exists(pptx_file):
        # Initialize PowerPoint application
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1  # Optional: Set to 1 to make PowerPoint visible
        
        # Open the PowerPoint file
        presentation = powerpoint.Presentations.Open(pptx_file)
        
        # Save as PDF (32 is the format for PDF)
        presentation.SaveAs(pdf_file, 32)
        presentation.Close()
        powerpoint.Quit()
        
        print(f"'{pptx_file}' has been converted to PDF!")
    else:
        print(f"File not found: {pptx_file}")

# Load Excel file with participant details
excel_path = "participation.xlsx"
data = pd.read_excel(excel_path)

# Load the PowerPoint template
template_path = "Certificate_participation.pptx"
output_folder = "Participation"
os.makedirs(output_folder, exist_ok=True)

# Define font and size settings for different tags
font_settings = {
    "<<Name>>": {"font": "Halant", "size": Pt(36)},
    "<<College>>": {"font": "Halant Bold", "size": Pt(20), "bold": True},
}

for _, row in data.iterrows():
    # Extract participant details
    name = row['Name']
    college = row['College']

    # Open the PowerPoint template
    prs = Presentation(template_path)

    # Replace placeholders in all slides and apply formatting
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:  # Check if the shape contains text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Check if the text in the run contains any of the placeholders
                        for placeholder, settings in font_settings.items():
                            if placeholder in run.text:
                                # Replace the placeholder with the actual value
                                run.text = run.text.replace(placeholder, locals()[placeholder[2:-2].lower()])
                                
                                # Apply the font and size to the run
                                run.font.name = settings["font"]
                                run.font.size = settings["size"]
                                
                                # Center align text if the placeholder is for College
                                if placeholder == "<<College>>":
                                    paragraph.alignment = 1  # 1 corresponds to center alignment

    # Save the personalized presentation
    pptx_file = os.path.join(os.getcwd(), output_folder, f"{name}_certificate.pptx")
    prs.save(pptx_file)

    pdf_file = os.path.join(os.getcwd(), output_folder, f"{name}_certificate.pdf")

    pptx_to_pdf(pptx_file, pdf_file)

print("Certificates generated successfully!")
print("Certificates generated and saved as PDF successfully!")
